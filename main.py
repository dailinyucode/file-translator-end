from dotenv import load_dotenv
load_dotenv()  # 添加在文件开头

from fastapi import FastAPI, UploadFile, File, Response, Form
from fastapi.middleware.cors import CORSMiddleware
from docx import Document
from pptx import Presentation
import PyPDF2
import io
import os
from tenacity import retry, stop_after_attempt, wait_exponential
import asyncio
from openai import AsyncOpenAI
from urllib.parse import quote  # 添加这个导入

app = FastAPI()

# 配置跨域
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # 允许所有来源，实际使用时建议指定具体域名
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 修改 OpenAI 客户端初始化
client = AsyncOpenAI(
    api_key="no-key",
    base_url=os.getenv('OPENAI_BASE_URL')  # 从环境变量获取 base_url
)

# @retry(
#     stop=stop_after_attempt(3),
#     wait=wait_exponential(multiplier=1, min=4, max=10)
# )
async def translate_text(text: str, target_lang: str) -> str:
    """使用 OpenAI 进行翻译"""
    try:
        prompt = f"请将以下文本翻译成{target_lang}，只返回翻译结果，不要包含其他内容：\n\n{text}"
        response = await client.chat.completions.create(
            model="internlm2.5-chat",
            messages=[
                {"role": "system", "content": "你是一个专业的翻译助手，请直接提供翻译结果，不要添加任何解释或额外内容。"},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            max_tokens=1500
        )
    
        translated_text = response.choices[0].message.content.strip()
        return translated_text
        
    except Exception as e:
        print(f"翻译过程发生错误: {str(e)}")
        raise Exception(f"翻译失败: {str(e)}")

async def process_docx(file: bytes, target_lang: str) -> bytes:
    """处理 Word 文档，保留原有格式"""
    doc = Document(io.BytesIO(file))
    
    try:
        # 批量收集需要翻译的文本
        texts_to_translate = []
        text_locations = []
        
        # 收集段落文本，同时保存段落格式信息
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                texts_to_translate.append(paragraph.text)
                # 保存段落的样式名称和格式
                text_locations.append({
                    'type': 'paragraph',
                    'element': paragraph,
                    'style': paragraph.style.name,
                    'runs': [(run.text, run.font.name, run.font.size, 
                             run.font.bold, run.font.italic, 
                             run.font.underline) for run in paragraph.runs]
                })
        
        # 收集表格文本
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts_to_translate.append(cell.text)
                        # 保存单元格中的格式信息
                        text_locations.append({
                            'type': 'cell',
                            'element': cell,
                            'runs': [(run.text, run.font.name, run.font.size,
                                     run.font.bold, run.font.italic,
                                     run.font.underline) 
                                    for paragraph in cell.paragraphs 
                                    for run in paragraph.runs]
                        })
        
        # 批量翻译
        for i, (text, location) in enumerate(zip(texts_to_translate, text_locations)):
            try:
                await asyncio.sleep(1)  # 添加延迟
                translated_text = await translate_text(text, target_lang)
                
                if location['type'] == 'paragraph':
                    paragraph = location['element']
                    paragraph.clear()
                    # 恢复段落样式
                    paragraph.style = doc.styles[location['style']]
                    
                    # 如果原文只有一个 run，直接添加翻译后的文本
                    if len(location['runs']) <= 1:
                        run = paragraph.add_run(translated_text)
                        if location['runs']:
                            # 应用原始格式
                            font = run.font
                            _, font_name, font_size, bold, italic, underline = location['runs'][0]
                            if font_name: font.name = font_name
                            if font_size: font.size = font_size
                            font.bold = bold
                            font.italic = italic
                            font.underline = underline
                    else:
                        # 如果原文有多个 run，尝试保持相似的分段
                        # 这里使用简单的按比例分配方式
                        total_len = sum(len(run[0]) for run in location['runs'])
                        if total_len > 0:
                            chars_per_run = [len(run[0]) / total_len for run in location['runs']]
                            start = 0
                            for j, (_, font_name, font_size, bold, italic, underline) in enumerate(location['runs']):
                                end = int(start + len(translated_text) * chars_per_run[j])
                                if j == len(location['runs']) - 1:
                                    end = len(translated_text)
                                run = paragraph.add_run(translated_text[start:end])
                                # 应用原始格式
                                font = run.font
                                if font_name: font.name = font_name
                                if font_size: font.size = font_size
                                font.bold = bold
                                font.italic = italic
                                font.underline = underline
                                start = end
                
                else:  # cell
                    cell = location['element']
                    cell.text = ""  # 清除原有内容
                    paragraph = cell.paragraphs[0]
                    run = paragraph.add_run(translated_text)
                    # 应用第一个 run 的格式（如果存在）
                    if location['runs']:
                        _, font_name, font_size, bold, italic, underline = location['runs'][0]
                        font = run.font
                        if font_name: font.name = font_name
                        if font_size: font.size = font_size
                        font.bold = bold
                        font.italic = italic
                        font.underline = underline
                
                if (i + 1) % 10 == 0:
                    await asyncio.sleep(1)
                    
            except Exception as e:
                print(f"翻译失败: {str(e)}")
                print(f"失败的文本: {text[:100]}...")
                continue
        
        # 保存翻译后的文档
        output = io.BytesIO()
        doc.save(output)
        return output.getvalue()
        
    except Exception as e:
        raise Exception(f"处理文档时出错: {str(e)}")

async def process_pptx(file: bytes, target_lang: str) -> bytes:
    """处理 PPT 文档"""
    prs = Presentation(io.BytesIO(file))
    
    # 遍历所有幻灯片
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                translated_text = await translate_text(shape.text, target_lang)
                shape.text = translated_text
    
    # 保存翻译后的文档
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()

async def process_pdf(file: bytes, target_lang: str) -> bytes:
    """处理 PDF 文档"""
    # 注意：PDF 的处理会比较复杂，这里只是一个简单示例
    reader = PyPDF2.PdfReader(io.BytesIO(file))
    writer = PyPDF2.PdfWriter()
    
    for page in reader.pages:
        text = page.extract_text()
        if text.strip():
            translated_text = await translate_text(text, target_lang)
            # 创建一个新页面with translated text
            # 注意：这种方式会丢失原始格式
            writer.add_page(page)
    
    output = io.BytesIO()
    writer.write(output)
    return output.getvalue()

@app.post("/translate-file")
async def translate_file(
    file: UploadFile = File(...),
    target_lang: str = Form("en")
):
    """文件翻译接口"""
    content = await file.read()
    file_extension = os.path.splitext(file.filename)[1].lower()
    
    print(f"接收到的目标语言: {target_lang}")
    
    try:
        if file_extension == '.docx':
            translated_content = await process_docx(content, target_lang)
            media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        elif file_extension == '.pptx':
            translated_content = await process_pptx(content, target_lang)
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        elif file_extension == '.pdf':
            translated_content = await process_pdf(content, target_lang)
            media_type = "application/pdf"
        elif file_extension == '.txt':
            # 修改这里的编码处理
            try:
                text = content.decode('utf-8')
            except UnicodeDecodeError:
                text = content.decode('gbk', errors='ignore')  # 尝试使用 GBK 编码
                
            translated_text = await translate_text(text, target_lang)
            translated_content = translated_text.encode('utf-8')
            media_type = "text/plain; charset=utf-8"
        else:
            return {"error": "不支持的文件格式"}
        
        # 修改文件名处理部分
        translated_filename = f"translated_{file.filename}"
        encoded_filename = quote(translated_filename)  # URL 编码文件名
        
        return Response(
            content=translated_content,
            media_type=media_type,
            headers={
                "Content-Disposition": f"attachment; filename={encoded_filename}"
            }
        )
    except Exception as e:
        print(f"错误详情: {str(e)}")
        return {"error": f"处理文件时出错: {str(e)}"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8989)
